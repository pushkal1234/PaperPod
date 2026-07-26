import logging
import os
import re
import uuid
from collections import Counter
from pathlib import Path

import PyPDF2
from docx import Document as DocxDocument
from docx.document import Document as _DocxDocumentClass
from docx.table import Table as _DocxTable
from docx.text.paragraph import Paragraph as _DocxParagraph

from app.config import settings

logger = logging.getLogger("paperpod")

# Minimum embedded-image dimension (px) to treat a PDF image as a real figure
# rather than a logo/icon — avoids sending trivial images to the vision model.
_MIN_FIGURE_DIM = 200
# Number of vector paths on a page that signals a drawn diagram/chart.
_VECTOR_DIAGRAM_THRESHOLD = 30
# Resolution to render figure pages at before sending to the vision model.
_FIGURE_RENDER_DPI = 150


# NUL and other C0 control chars that must never reach the DB or TTS. Postgres
# text columns reject 0x00 outright ("invalid byte sequence for encoding UTF8:
# 0x00"), and PyPDF2 / OCR routinely emit them for scanned or oddly-encoded
# PDFs. We keep \t (0x09), \n (0x0a), and \r (0x0d); everything else in the C0
# range plus the NUL byte is stripped.
_CONTROL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def clean_extracted_text(text: str) -> str:
    """Strip NUL/control bytes that break Postgres storage and downstream calls.

    Applied to every source of ``raw_text`` (PDF/DOCX/TXT extraction, pasted
    text, and image OCR) so a stray 0x00 can never crash the storage step.
    """
    if not text:
        return text
    return _CONTROL_CHARS_RE.sub("", text)


# Standalone page-number / running-folio lines: "12", "- 12 -", "Page 12",
# "Page 12 of 34", "12 / 34". These carry no semantic value in a podcast.
_PAGE_NUM_RE = re.compile(
    r"^\s*(?:[-–—]\s*)?(?:page\s+)?\d{1,4}(?:\s*/\s*\d{1,4}|\s+of\s+\d{1,4})?\s*(?:[-–—])?\s*$",
    re.IGNORECASE,
)
# A word hyphenated across a line break in a PDF: "informa-\ntion" -> "information".
# Includes the soft-hyphen (U+00AD) PDFs sometimes emit.
_HYPHEN_BREAK_RE = re.compile(r"(\w)[-\u00ad]\n(\w)")
# 3+ consecutive newlines collapse to exactly one blank line.
_MULTI_BLANK_RE = re.compile(r"\n{3,}")
# Below this size a doc has no meaningful boilerplate to reclaim and is cheap
# already — skip compaction entirely to keep short docs 100% untouched.
_COMPACTION_MIN_CHARS = 2000
# If compaction removes more than this fraction, treat it as an anomaly (e.g. a
# doc that is genuinely mostly repeated lines) and keep the original to be safe.
_COMPACTION_MAX_REMOVAL = 0.40


def compact_document_text(text: str) -> str:
    """Conservatively remove NON-SEMANTIC extraction noise before the LLM call.

    This is a *lossless-of-meaning* reduction — it only removes chrome that a
    human reader ignores and that wastes input tokens:
      - repeated running headers/footers (identical short lines on many pages)
      - standalone page-number / folio lines
      - hyphenation splits introduced by PDF line wrapping
      - trailing per-line whitespace and runs of blank lines

    It NEVER rewrites, summarizes, or drops prose. On any anomaly (empty result,
    output larger than input, or >40% removed) it returns the ORIGINAL text so a
    surprising input can never degrade the podcast. Disable via DOC_COMPACTION=0.
    """
    if not text or len(text) < _COMPACTION_MIN_CHARS:
        return text
    try:
        original = text
        # 1. Rejoin words split across a line break (lossless improvement — the
        #    tokenizer would otherwise see two fragments instead of one word).
        text = _HYPHEN_BREAK_RE.sub(r"\1\2", text)

        raw_lines = text.split("\n")

        # 2. Detect repeated running headers/footers: a short, non-sentence line
        #    that recurs many times is almost always page chrome (title bar,
        #    footer, section running-head), not content. Require >=4 repeats,
        #    <=70 chars, at least one alphanumeric char, and NOT ending in
        #    sentence punctuation (so real repeated sentences are never caught).
        norm = [ln.strip() for ln in raw_lines]
        counts = Counter(l for l in norm if l)
        repeated_chrome = {
            l
            for l, c in counts.items()
            if c >= 4 and len(l) <= 70 and l[-1] not in ".:!?" and any(ch.isalnum() for ch in l)
        }

        out: list[str] = []
        seen_chrome: set[str] = set()
        for ln in raw_lines:
            s = ln.rstrip()
            stripped = s.strip()
            if stripped and _PAGE_NUM_RE.match(stripped):
                continue  # drop standalone page-number lines
            if stripped in repeated_chrome:
                # Keep the FIRST occurrence (it may be the document title); drop
                # every later repeat.
                if stripped in seen_chrome:
                    continue
                seen_chrome.add(stripped)
            out.append(s)

        text = _MULTI_BLANK_RE.sub("\n\n", "\n".join(out)).strip()

        # Safety gates — never let compaction surprise the pipeline.
        if not text or len(text) > len(original):
            return original
        if len(text) < len(original) * (1 - _COMPACTION_MAX_REMOVAL):
            logger.warning(
                "[doc] compaction removed >%.0f%% (%d -> %d chars) — anomaly, keeping original",
                _COMPACTION_MAX_REMOVAL * 100,
                len(original),
                len(text),
            )
            return original
        return text
    except Exception as e:  # never break generation over a cleanup step
        logger.warning("[doc] compaction failed (%s) — using original text", e)
        return text


def extract_text(file_path: str, content_type: str) -> str:
    """Extract text from PDF, DOCX, or TXT files."""
    if content_type == "application/pdf" or file_path.endswith(".pdf"):
        text = _extract_pdf(file_path)
    elif content_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ) or file_path.endswith(".docx"):
        text = _extract_docx(file_path)
    elif content_type == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ) or file_path.endswith(".pptx"):
        text = _extract_pptx(file_path)
    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    return clean_extracted_text(text)


def _extract_pdf(file_path: str) -> str:
    text_parts = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        # Many PDFs are encrypted with an owner/permissions password but an EMPTY
        # user password — they open fine in any viewer. Decrypt with "" so we can
        # read them; if that fails the PDF is genuinely password-protected.
        if reader.is_encrypted:
            try:
                if reader.decrypt("") == PyPDF2.PasswordType.NOT_DECRYPTED:
                    raise RuntimeError(
                        "This PDF is password-protected. Please remove the password "
                        "(open it and re-save/print to PDF) and upload it again."
                    )
            except RuntimeError:
                raise
            except Exception as e:  # noqa: BLE001 — surface a clear, user-facing message
                logger.warning(f"[PDF] Could not decrypt encrypted PDF: {e}")
                raise RuntimeError(
                    "This PDF is encrypted and could not be opened. Please remove "
                    "the password (open it and re-save/print to PDF) and try again."
                )
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    text = "\n\n".join(text_parts)

    # Enrich with spoken-language descriptions of diagrams/charts/figures that
    # PyPDF2 (text layer only) cannot read, so the podcast can narrate visuals.
    visuals = _describe_pdf_visuals(file_path)
    if visuals:
        text = f"{text}\n\n## Visual elements (diagrams, charts, and figures)\n{visuals}"
    return text


def _describe_pdf_visuals(file_path: str) -> str:
    """Render figure-bearing PDF pages and get Gemini descriptions of them.

    Best-effort: returns "" (never raises) if vision is disabled/unconfigured,
    PyMuPDF is unavailable, the PDF has no real figures, or the call fails.
    """
    if not settings.PDF_VISION_EXTRACTION or not settings.GOOGLE_API_KEY:
        return ""
    try:
        import fitz  # PyMuPDF — imported lazily so it's only needed for PDFs
    except ImportError:
        logger.warning("[Vision] PyMuPDF not installed; skipping figure descriptions")
        return ""

    try:
        pages = _render_figure_pages(fitz, file_path)
    except Exception as e:  # noqa: BLE001 — best-effort enhancement
        logger.warning(f"[Vision] Could not scan PDF for figures ({e})")
        return ""

    if not pages:
        return ""

    from app.services.image_service import describe_pdf_figures
    return describe_pdf_figures(pages)


def _page_has_visual(page) -> bool:
    """True if a PDF page likely contains a real figure/diagram/chart.

    Catches (a) sizeable embedded raster images (ignoring tiny logos/icons) and
    (b) vector-drawn diagrams with many paths. A full-page scan also qualifies.
    """
    for img in page.get_images(full=True):
        width, height = img[2], img[3]
        if width >= _MIN_FIGURE_DIM and height >= _MIN_FIGURE_DIM:
            return True
    try:
        if len(page.get_drawings()) >= _VECTOR_DIAGRAM_THRESHOLD:
            return True
    except Exception:  # noqa: BLE001 — drawing extraction is optional
        pass
    return False


def _render_figure_pages(fitz, file_path: str) -> list[tuple[int, bytes]]:
    """Return [(page_number, png_bytes)] for figure-bearing pages, capped.

    Skips very long PDFs entirely (latency/cost guard).
    """
    rendered: list[tuple[int, bytes]] = []
    doc = fitz.open(file_path)
    try:
        if doc.page_count > settings.PDF_VISION_MAX_PAGES:
            logger.info(
                f"[Vision] Skipping figure scan: {doc.page_count} pages "
                f"> PDF_VISION_MAX_PAGES ({settings.PDF_VISION_MAX_PAGES})"
            )
            return []
        matrix = fitz.Matrix(_FIGURE_RENDER_DPI / 72.0, _FIGURE_RENDER_DPI / 72.0)
        for i in range(doc.page_count):
            if len(rendered) >= settings.PDF_VISION_MAX_FIGURES:
                break
            page = doc[i]
            if _page_has_visual(page):
                pix = page.get_pixmap(matrix=matrix)
                rendered.append((i + 1, pix.tobytes("png")))
        if rendered:
            logger.info(f"[Vision] Found {len(rendered)} figure page(s) to describe")
    finally:
        doc.close()
    return rendered


def _iter_block_items(parent):
    """Yield paragraphs and tables from a docx body in document order.

    python-docx exposes paragraphs and tables in separate collections that do
    not preserve their relative order. Walking the underlying XML body lets us
    reconstruct the original sequence so table content (e.g. procedure steps)
    is not dropped.
    """
    if isinstance(parent, _DocxDocumentClass):
        parent_elm = parent.element.body
    else:
        parent_elm = parent._element

    for child in parent_elm.iterchildren():
        if child.tag.endswith("}p"):
            yield _DocxParagraph(child, parent)
        elif child.tag.endswith("}tbl"):
            yield _DocxTable(child, parent)


def _dedupe_cell_text(raw: str) -> str:
    """Collapse repeated lines within a cell (vertical-merge artifact).

    python-docx returns a vertically merged cell's text once per row it spans,
    e.g. "PANEL ENGINEER\nPANEL ENGINEER". Keep consecutive unique lines only,
    then flatten internal whitespace.
    """
    seen_lines = []
    prev = None
    for line in raw.splitlines():
        cleaned = " ".join(line.split())
        if not cleaned or cleaned == prev:
            continue
        seen_lines.append(cleaned)
        prev = cleaned
    return " ".join(seen_lines)


def _render_table(table: _DocxTable) -> str:
    """Render a docx table as pipe-delimited rows the LLM can read.

    Collapses duplicate cell values from merged cells (column spans) and
    repeated lines within a cell (row spans).
    """
    rows = []
    for row in table.rows:
        cells = []
        prev = None
        for cell in row.cells:
            text = _dedupe_cell_text(cell.text)
            if text and text == prev:
                continue
            cells.append(text)
            prev = text
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _extract_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    parts = []
    for block in _iter_block_items(doc):
        if isinstance(block, _DocxParagraph):
            if block.text.strip():
                parts.append(block.text.strip())
        elif isinstance(block, _DocxTable):
            rendered = _render_table(block)
            if rendered.strip():
                parts.append(rendered)
    return "\n\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    # Lazy import so python-pptx only loads when a deck is actually uploaded.
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    parts = []
    images = []  # (slide_no, blob, content_type) for the vision pass

    for i, slide in enumerate(prs.slides, start=1):
        chunk = []
        for shape in _walk_shapes(slide.shapes, MSO_SHAPE_TYPE):
            if getattr(shape, "has_table", False):
                chunk.append(_pptx_table(shape.table))
            elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                chunk.append(shape.text_frame.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    images.append((i, shape.image.blob, shape.image.content_type))
                except Exception:
                    pass
        # Speaker notes are often where the real explanation lives.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunk.append(f"Notes: {notes}")
        if chunk:
            parts.append(f"Slide {i}:\n" + "\n".join(chunk))

    text = "\n\n".join(parts)

    # Same Gemini vision path as PDFs; pictures come straight from the file so
    # no page rendering is needed.
    visuals = _describe_pptx_images(images)
    if visuals:
        text = f"{text}\n\n## Visual elements (diagrams, charts, and figures)\n{visuals}"
    return text


def _walk_shapes(shapes, mso_type):
    # Flatten groups so nested text and pictures aren't missed.
    for shape in shapes:
        if shape.shape_type == mso_type.GROUP:
            yield from _walk_shapes(shape.shapes, mso_type)
        else:
            yield shape


def _pptx_table(table) -> str:
    rows = []
    for row in table.rows:
        cells = [c.text.strip() for c in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _describe_pptx_images(images: list[tuple[int, bytes, str]]) -> str:
    """Describe raster pictures embedded in slides via Gemini (best-effort)."""
    if not settings.PDF_VISION_EXTRACTION or not settings.GOOGLE_API_KEY or not images:
        return ""
    from io import BytesIO
    from PIL import Image

    pages = []
    for slide_no, blob, content_type in images:
        if len(pages) >= settings.PDF_VISION_MAX_FIGURES:
            break
        if "emf" in content_type or "wmf" in content_type:
            continue  # vector art Gemini can't read
        try:
            img = Image.open(BytesIO(blob))
            if min(img.size) < _MIN_FIGURE_DIM:
                continue  # skip logos/icons/bullets
            img = img.convert("RGB")
            if max(img.size) > 1280:
                r = 1280 / max(img.size)
                img = img.resize((int(img.width * r), int(img.height * r)))
            buf = BytesIO()
            img.save(buf, format="PNG")
            pages.append((slide_no, buf.getvalue()))
        except Exception:
            continue
    if not pages:
        return ""
    from app.services.image_service import describe_pdf_figures
    return describe_pdf_figures(pages)


def chunk_text(text: str, chunk_size: int = 1500, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks for embedding."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap
    return chunks


def save_upload(file_bytes: bytes, filename: str) -> str:
    """Save uploaded file to disk, return path."""
    ext = Path(filename).suffix
    unique_name = f"{uuid.uuid4()}{ext}"
    file_path = os.path.join(settings.UPLOAD_DIR, unique_name)
    with open(file_path, "wb") as f:
        f.write(file_bytes)
    return file_path
